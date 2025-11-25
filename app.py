# app.py - CORRECTED ADAPTIVE LOGIC
import os
import json
import sqlite3
import time
import random
from datetime import datetime, timedelta
from urllib.parse import urlparse
from functools import wraps
import requests
from bs4 import BeautifulSoup
import PyPDF2
from pptx import Presentation

from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename

from simple_adaptive_engine import SimpleAdaptiveEngine 
from gemini_engine import GeminiQuizEngine

# --- Configuration & Initialization ---
app = Flask(__name__)
app.config['SECRET_KEY'] = 'smart-quizzer-final-key-45678'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
app.config['DATABASE'] = 'quizzes.db'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

adaptive_engine = SimpleAdaptiveEngine()
quiz_engine = GeminiQuizEngine()

# Motivational quotes for enhanced user experience
MOTIVATIONAL_QUOTES = [
    "Every expert was once a beginner. Keep learning! 🚀",
    "Knowledge is power! Your progress is impressive! 💪",
    "Great job! Every quiz makes you smarter! 🧠",
    "Learning is a journey, not a destination. Well done! 🌟",
    "You're building knowledge brick by brick! Amazing work! 🏗️",
    "Success is the sum of small efforts repeated daily! 🔥",
    "Your brain is getting stronger with every question! 💡",
    "Learning never exhausts the mind - Leonardo da Vinci 🎨",
    "The more you learn, the more you earn! Future you says thanks! 💰",
    "You're not just taking quizzes, you're building your future! 🌈",
    "Mistakes are proof that you are trying! Keep going! 💫",
    "Your potential is endless! Believe in yourself! 🌠",
    "Every question answered is a step toward mastery! 📚",
    "You're capable of amazing things! 🌟",
    "Progress, not perfection! You're doing great! 🎯"
]

# --- Utility Functions ---

def is_url(content):
    """Checks if a string is a valid URL."""
    try:
        result = urlparse(content)
        return all([result.scheme in ['http', 'https'], result.netloc])
    except ValueError:
        return False

def extract_text_from_file(filepath):
    """
    Extracts real text content from PDF, PPTX, and TXT files.
    """
    filename = os.path.basename(filepath)
    ext = filename.split('.')[-1].lower()
    text_content = ""
    
    try:
        # 1. Handle PDF Files
        if ext == 'pdf':
            print(f"DEBUG: Processing PDF - {filename}")
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                # Limit to first 20 pages to avoid overloading Gemini
                for page in reader.pages[:20]:
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
        
        # 2. Handle PowerPoint (PPTX) Files
        elif ext == 'pptx':
            print(f"DEBUG: Processing PPTX - {filename}")
            prs = Presentation(filepath)
            # Limit to first 30 slides
            for slide in prs.slides[:30]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

        # 3. Handle Text (TXT) Files
        elif ext == 'txt':
            print(f"DEBUG: Processing TXT - {filename}")
            with open(filepath, 'r', encoding='utf-8') as f:
                text_content = f.read()

        # 4. Handle Unsupported Files
        else:
            return f"Error: Unsupported file type '{ext}'. Please upload PDF, PPTX, or TXT."

        # Final cleanup: Limit to ~10,000 characters for Gemini API safety
        cleaned_text = text_content.strip()
        if not cleaned_text:
            return "Error: The file appears to be empty or contains only images (no readable text)."
            
        return cleaned_text[:10000]

    except Exception as e:
        print(f"Error extracting content: {e}")
        return f"Error processing file: {str(e)}"

def get_db_connection(dict_cursor=True):
    """Get database connection with optional dictionary cursor."""
    conn = sqlite3.connect(app.config['DATABASE'])
    if dict_cursor:
        conn.row_factory = sqlite3.Row
    return conn

def init_db():
    """Initialize database with all required tables and migrations."""
    # 1. Open Connection
    conn = get_db_connection(dict_cursor=False)
    c = conn.cursor()
    
    try:
        # 2. Create Tables
        c.executescript('''
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT UNIQUE NOT NULL,
                email TEXT UNIQUE NOT NULL,
                phone TEXT,
                password TEXT NOT NULL,
                security_q TEXT,
                security_a TEXT,
                skill_level TEXT DEFAULT 'Medium',
                status TEXT DEFAULT 'active',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            
            CREATE TABLE IF NOT EXISTS quizzes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER NOT NULL,
                title TEXT NOT NULL,
                topic TEXT NOT NULL,
                content TEXT NOT NULL,
                questions TEXT NOT NULL,
                difficulty TEXT DEFAULT 'Medium',
                quiz_type TEXT DEFAULT 'adaptive',
                score REAL DEFAULT 0,
                status TEXT DEFAULT 'in_progress',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users (id)
            );
            
            CREATE TABLE IF NOT EXISTS performances (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER NOT NULL,
                topic TEXT NOT NULL,
                difficulty TEXT NOT NULL,
                accuracy REAL DEFAULT 0.0,
                total_questions INTEGER DEFAULT 0,
                correct_answers INTEGER DEFAULT 0,
                average_response_time REAL DEFAULT 0.0,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            
            CREATE TABLE IF NOT EXISTS admins (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT UNIQUE NOT NULL,
                password TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            
            CREATE TABLE IF NOT EXISTS question_feedback (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                quiz_id INTEGER,
                question_id TEXT,
                user_id INTEGER,
                feedback_type TEXT,
                comment TEXT,
                flagged INTEGER DEFAULT 0,
                resolved INTEGER DEFAULT 0,
                question_text TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (quiz_id) REFERENCES quizzes (id),
                FOREIGN KEY (user_id) REFERENCES users (id)
            );
        ''')
        
        # 3. Helper to safely add columns (Migrations)
        def add_column_if_not_exists(table, column, definition):
            try:
                c.execute(f"SELECT {column} FROM {table} LIMIT 1")
            except sqlite3.OperationalError:
                print(f"🔄 Migrating DB: Adding {column} to {table}...")
                c.execute(f'ALTER TABLE {table} ADD COLUMN {column} {definition}')

        # Run Migrations (Ensure columns exist)
        add_column_if_not_exists('users', 'phone', 'TEXT')
        add_column_if_not_exists('users', 'security_q', 'TEXT')
        add_column_if_not_exists('users', 'security_a', 'TEXT')
        add_column_if_not_exists('users', 'skill_level', 'TEXT DEFAULT "Medium"')
        add_column_if_not_exists('users', 'status', 'TEXT DEFAULT "active"')
        
        # 4. Create Default Admin (ONLY if not exists)
        default_admin_username = 'admin'
        default_admin_password = 'Admin@123'

        try:
            c.execute('SELECT id FROM admins WHERE username = ?', (default_admin_username,))
            existing_admin = c.fetchone()
            
            if not existing_admin:
                c.execute(
                    'INSERT INTO admins (username, password) VALUES (?, ?)',
                    (default_admin_username, generate_password_hash(default_admin_password))
                )
                print("✅ Default Admin created.")
        except Exception as e:
            print(f"⚠️ Admin check warning: {e}")

        # 5. Commit Changes
        conn.commit()

    except Exception as e:
        print(f"❌ Database Initialization Error: {e}")
        
    finally:
        # 6. Close Connection (ONLY HERE at the very end)
        conn.close()

def is_logged_in():
    """Check if user is logged in."""
    return 'user_id' in session

def is_admin_logged_in():
    """Check if admin is logged in."""
    return session.get('is_admin')

def get_current_user_data(user_id):
    """Get current user data from database."""
    with get_db_connection() as conn:
        user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    return user

def get_user_skill_level(user_id):
    """Get user's current skill level."""
    user_data = get_current_user_data(user_id)
    return user_data['skill_level'] if user_data else 'Medium'

def evaluate_answer(question, user_answer):
    """Enhanced answer evaluation for multiple question types."""
    question_type = question.get('question_type', 'mcq')
    correct_answer = question.get('correct_answer', '')
    
    if not user_answer or not correct_answer:
        return False
        
    user_answer = user_answer.strip()
    correct_answer = correct_answer.strip()
    
    if question_type == 'checkbox':
        # For checkbox questions, user_answer is comma-separated string of selected options
        user_answers = [ans.strip().lower() for ans in user_answer.split(',')] if user_answer else []
        correct_answers = [ans.strip().lower() for ans in correct_answer.split(',')] if correct_answer else []
        
        # Check if all correct answers are selected and no incorrect ones
        user_set = set(user_answers)
        correct_set = set(correct_answers)
        
        return user_set == correct_set
        
    elif question_type == 'true_false':
        return user_answer.lower() == correct_answer.lower()
        
    elif question_type in ['mcq', 'dropdown']:
        return user_answer.lower() == correct_answer.lower()
        
    elif question_type == 'short_answer':
        # For short answer, do case-insensitive partial matching
        return user_answer.lower() in correct_answer.lower()
    
    return False

def generate_ai_suggestions(score, weak_topics, difficulty_performance, quiz_topic):
    """Generate AI-powered learning suggestions."""
    suggestions = []
    
    # Performance-based suggestions
    if score >= 90:
        suggestions.append({
            'type': 'challenge',
            'title': 'Ready for Advanced Challenges!',
            'description': f'You\'ve mastered the basics of {quiz_topic}. Try creating quizzes with "Hard" difficulty or explore related advanced topics.',
            'icon': '🚀',
            'priority': 1
        })
    elif score >= 75:
        suggestions.append({
            'type': 'practice',
            'title': 'Solid Foundation - Time to Level Up!',
            'description': f'You have a good grasp of {quiz_topic}. Focus on medium-difficulty questions and review the topics you missed.',
            'icon': '💪',
            'priority': 1
        })
    else:
        suggestions.append({
            'type': 'foundation',
            'title': 'Build Strong Foundations',
            'description': f'Focus on understanding the core concepts of {quiz_topic}. Start with easier questions and gradually increase difficulty.',
            'icon': '🏗️',
            'priority': 1
        })
    
    # Weak topics suggestions
    if weak_topics:
        top_weak_topic = max(weak_topics, key=weak_topics.get)
        suggestions.append({
            'type': 'focus',
            'title': f'Focus on: {top_weak_topic}',
            'description': f'You missed {weak_topics[top_weak_topic]} questions on {top_weak_topic}. Consider reviewing this area specifically.',
            'icon': '🎯',
            'priority': 2
        })
    
    # Difficulty-based suggestions
    for difficulty, stats in difficulty_performance.items():
        accuracy = (stats['correct'] / stats['total']) * 100 if stats['total'] > 0 else 0
        if accuracy < 70:
            suggestions.append({
                'type': 'difficulty',
                'title': f'Improve {difficulty} Level Questions',
                'description': f'Your accuracy on {difficulty.lower()} questions is {accuracy:.1f}%. Practice more {difficulty.lower()} level material.',
                'icon': '📊',
                'priority': 2
            })
    
    # General learning strategies
    learning_strategies = [
        {
            'type': 'strategy',
            'title': 'Spaced Repetition',
            'description': 'Review the topics you missed in 24 hours, then again in 3 days for better retention.',
            'icon': '⏰',
            'priority': 3
        },
        {
            'type': 'strategy', 
            'title': 'Active Recall',
            'description': 'Try to explain the concepts you learned to someone else - it strengthens memory.',
            'icon': '🗣️',
            'priority': 3
        },
        {
            'type': 'strategy',
            'title': 'Interleaved Practice',
            'description': 'Mix different topics in your study sessions for better long-term retention.',
            'icon': '🔄',
            'priority': 3
        }
    ]
    
    suggestions.extend(learning_strategies[:2])  # Add 2 strategies
    
    # Next quiz suggestion
    next_difficulty = "Medium"
    if score >= 80:
        next_difficulty = "Hard"
    elif score <= 50:
        next_difficulty = "Easy"
    
    suggestions.append({
        'type': 'next_step',
        'title': 'Recommended Next Quiz',
        'description': f'Try another {quiz_topic} quiz with {next_difficulty} difficulty to continue your learning journey.',
        'icon': '➡️',
        'priority': 1
    })
    
    # Sort by priority
    suggestions.sort(key=lambda x: x['priority'])
    
    return suggestions

# --- Admin Authentication Middleware ---
def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not is_admin_logged_in():
            return redirect(url_for('admin_login'))
        return f(*args, **kwargs)
    return decorated_function

def scrape_url_content(url):
    """Fetches text content from a URL."""
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Remove scripts and styles
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
            
        # Get text
        text = ' '.join([p.get_text() for p in soup.find_all('p')])
        
        # Limit text to 8000 characters to fit Gemini's context window
        return text[:8000]
    except Exception as e:
        print(f"Scraping Error: {e}")
        return None

# --- ROUTES ---

@app.route('/')
def index():
    """Home page route."""
    if is_logged_in():
        return redirect('/dashboard')
    return render_template('index.html')

@app.route('/admin/login', methods=['GET', 'POST'])
def admin_login():
    """Admin login route."""
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '').strip()

        with get_db_connection() as conn:
            admin = conn.execute('SELECT * FROM admins WHERE username = ?', (username,)).fetchone()

        if admin and check_password_hash(admin['password'], password):
            session.clear()
            session['is_admin'] = True
            session['admin_id'] = admin['id']
            session['admin_username'] = admin['username']
            flash('Admin login successful.', 'success')
            return redirect(url_for('admin_dashboard'))

        flash('Invalid admin credentials.', 'danger')

    return render_template('admin_login.html')

@app.route('/admin/logout')
def admin_logout():
    """Admin logout route."""
    session.clear()
    flash('Admin logged out successfully.', 'info')
    return redirect(url_for('admin_login'))

@app.route('/admin/profile')
def admin_profile():
    if not session.get('admin_logged_in'):
        return redirect(url_for('admin_login'))

    # if you want admin details, you can send them like this:
    admin_name = session.get('admin_username', 'Admin')

    return render_template('admin_profile.html', admin_name=admin_name)


@app.route('/admin/dashboard')
@admin_required
def admin_dashboard():
    """Enhanced admin dashboard with dynamic data."""
    if not is_admin_logged_in():
        return redirect(url_for('admin_login'))

    with get_db_connection() as conn:
        # Platform statistics
        total_users = conn.execute('SELECT COUNT(*) FROM users WHERE username != ?', ('admin',)).fetchone()[0]
        
        # Get active users (last 30 days) - users who completed quizzes
        thirty_days_ago = (datetime.now() - timedelta(days=30)).strftime('%Y-%m-%d %H:%M:%S')
        active_users = conn.execute(
            "SELECT COUNT(DISTINCT user_id) FROM quizzes WHERE created_at >= ? AND status = 'completed'",
            (thirty_days_ago,)
        ).fetchone()[0]
        
        total_quizzes = conn.execute('SELECT COUNT(*) FROM quizzes').fetchone()[0]
        completed_quizzes = conn.execute(
            "SELECT COUNT(*) FROM quizzes WHERE status = 'completed'"
        ).fetchone()[0]
        
        total_feedback = conn.execute('SELECT COUNT(*) FROM question_feedback').fetchone()[0]
        flagged_open = conn.execute(
            'SELECT COUNT(*) FROM question_feedback WHERE flagged = 1 AND resolved = 0'
        ).fetchone()[0]

        # Feedback analysis
        feedback_breakdown = conn.execute(
            'SELECT feedback_type, COUNT(*) as count FROM question_feedback GROUP BY feedback_type'
        ).fetchall()

        # User management data (excluding admin)
        user_rows = conn.execute(
            'SELECT id, username, email, skill_level, status, created_at FROM users WHERE username != ? ORDER BY created_at DESC LIMIT 10',
            ('admin',)
        ).fetchall()

        # Recent quiz content
        content_rows = conn.execute(
            'SELECT id, title, topic, status, created_at FROM quizzes ORDER BY created_at DESC LIMIT 5'
        ).fetchall()

        # Flagged questions
        flagged_rows = conn.execute(
            '''
            SELECT qf.*, u.username AS reported_by, q.title AS quiz_title
            FROM question_feedback qf
            LEFT JOIN users u ON qf.user_id = u.id
            LEFT JOIN quizzes q ON qf.quiz_id = q.id
            WHERE qf.flagged = 1 AND qf.resolved = 0
            ORDER BY qf.created_at DESC
            '''
        ).fetchall()

        # Leaderboard
        leaderboard_rows = conn.execute(
            '''
            SELECT u.id,
                   u.username,
                   COUNT(q.id) AS completed_quizzes,
                   AVG(q.score) AS avg_score,
                   MAX(q.score) AS best_score
            FROM users u
            JOIN quizzes q ON q.user_id = u.id AND q.status = 'completed'
            WHERE u.username != 'admin'
            GROUP BY u.id
            ORDER BY avg_score DESC, completed_quizzes DESC, best_score DESC
            LIMIT 5
            '''
        ).fetchall()

    # Process feedback sentiment
    feedback_counts = {'positive': 0, 'negative': 0, 'neutral': 0}
    for row in feedback_breakdown:
        f_type = (row['feedback_type'] or '').lower()
        if f_type in feedback_counts:
            feedback_counts[f_type] += row['count']

    sentiment_total = feedback_counts['positive'] + feedback_counts['negative']
    positive_pct = (feedback_counts['positive'] / sentiment_total * 100) if sentiment_total else 0
    negative_pct = (feedback_counts['negative'] / sentiment_total * 100) if sentiment_total else 0

    stats = {
        'total_users': total_users,
        'active_users': active_users,
        'total_quizzes': total_quizzes,
        'completed_quizzes': completed_quizzes,
        'total_feedback': total_feedback,
        'flagged_open': flagged_open,
        'positive_pct': round(positive_pct, 1),
        'negative_pct': round(negative_pct, 1)
    }

    # Convert to dictionaries for template
    users = [dict(row) for row in user_rows]
    content_items = [dict(row) for row in content_rows]
    flagged_items = [dict(row) for row in flagged_rows]
    leaderboard = [
        {
            'username': row['username'],
            'completed_quizzes': row['completed_quizzes'],
            'avg_score': round(row['avg_score'], 1) if row['avg_score'] is not None else 0,
        }
        for row in leaderboard_rows
    ]

    return render_template(
        'admin_dashboard.html',
        stats=stats,
        users=users,
        content_items=content_items,
        flagged_items=flagged_items,
        leaderboard=leaderboard
    )

@app.route('/admin/api/stats')
@admin_required
def get_live_stats():
    """API endpoint for live statistics"""
    conn = get_db_connection()
    
    try:
        # Get real-time stats
        total_users = conn.execute('SELECT COUNT(*) FROM users WHERE username != ?', ('admin',)).fetchone()[0]
        total_quizzes = conn.execute('SELECT COUNT(*) FROM quizzes').fetchone()[0]
        completed_quizzes = conn.execute(
            "SELECT COUNT(*) FROM quizzes WHERE status = 'completed'"
        ).fetchone()[0]
        flagged_open = conn.execute(
            'SELECT COUNT(*) FROM question_feedback WHERE flagged = 1 AND resolved = 0'
        ).fetchone()[0]
        
        # Get active users (last 30 days)
        thirty_days_ago = (datetime.now() - timedelta(days=30)).strftime('%Y-%m-%d %H:%M:%S')
        active_users = conn.execute(
            "SELECT COUNT(DISTINCT user_id) FROM quizzes WHERE created_at >= ? AND status = 'completed'",
            (thirty_days_ago,)
        ).fetchone()[0]
        
    except Exception as e:
        print(f"Error getting live stats: {e}")
        total_users = total_quizzes = completed_quizzes = flagged_open = active_users = 0
    
    finally:
        conn.close()
    
    return jsonify({
        'total_users': total_users,
        'total_quizzes': total_quizzes,
        'completed_quizzes': completed_quizzes,
        'flagged_open': flagged_open,
        'active_users': active_users,
        'timestamp': datetime.now().isoformat()
    })

@app.route('/admin/user/<int:user_id>')
@admin_required
def get_user_details(user_id):
    """Get detailed user information"""
    conn = get_db_connection()
    
    try:
        # Get user basic info
        user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
        if not user:
            return jsonify({'success': False, 'error': 'User not found'})
        
        # Get user's quiz attempts with details
        user_quizzes = conn.execute('''
            SELECT q.title, q.topic, q.difficulty, q.quiz_type,
                   q.score, q.created_at, q.status
            FROM quizzes q
            WHERE q.user_id = ?
            ORDER BY q.created_at DESC
        ''', (user_id,)).fetchall()
        
        # Get user statistics
        user_stats = conn.execute('''
            SELECT 
                COUNT(*) as total_attempts,
                AVG(score) as avg_score,
                MAX(score) as best_score,
                MIN(score) as worst_score
            FROM quizzes
            WHERE user_id = ? AND status = 'completed'
        ''', (user_id,)).fetchone()
        
        # Get user's feedback
        user_feedback = conn.execute('''
            SELECT feedback_type, comment, created_at
            FROM question_feedback
            WHERE user_id = ?
            ORDER BY created_at DESC
            LIMIT 5
        ''', (user_id,)).fetchall()
        
        return jsonify({
            'success': True,
            'user': dict(user),
            'quizzes': [dict(quiz) for quiz in user_quizzes],
            'stats': dict(user_stats) if user_stats else {},
            'feedback': [dict(fb) for fb in user_feedback]
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()

@app.route('/admin/resolve_flagged/<int:feedback_id>', methods=['POST'])
@admin_required
def admin_resolve_feedback(feedback_id):
    """Resolve flagged feedback."""
    if not is_admin_logged_in():
        return redirect(url_for('admin_login'))

    with get_db_connection(dict_cursor=False) as conn:
        conn.execute('UPDATE question_feedback SET resolved = 1 WHERE id = ?', (feedback_id,))
        conn.commit()

    flash('Feedback marked as resolved.', 'success')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/update_user/<int:user_id>', methods=['POST'])
@admin_required
def update_user(user_id):
    """Update user information"""
    data = request.json
    conn = get_db_connection()
    
    try:
        conn.execute('''
            UPDATE users 
            SET username = ?, email = ?, skill_level = ?, status = ?
            WHERE id = ?
        ''', (data.get('username'), data.get('email'), data.get('skill_level'), data.get('status'), user_id))
        
        conn.commit()
        return jsonify({'success': True, 'message': 'User updated successfully'})
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})
    finally:
        conn.close()

@app.route('/admin/delete_user/<int:user_id>', methods=['POST'])
@admin_required
def admin_delete_user(user_id):
    """Delete user account."""
    if not is_admin_logged_in():
        return redirect(url_for('admin_login'))

    with get_db_connection(dict_cursor=False) as conn:
        # Check if user exists and is not admin
        user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
        if not user:
            return jsonify({'success': False, 'error': 'User not found'})
        
        if user['username'] == 'admin':
            return jsonify({'success': False, 'error': 'Cannot delete admin user'})
        
        # Delete user's data
        conn.execute('DELETE FROM question_feedback WHERE user_id = ?', (user_id,))
        conn.execute('DELETE FROM performances WHERE user_id = ?', (user_id,))
        conn.execute('DELETE FROM quizzes WHERE user_id = ?', (user_id,))
        conn.execute('DELETE FROM users WHERE id = ?', (user_id,))
        conn.commit()

    return jsonify({'success': True, 'message': 'User deleted successfully'})

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        phone = request.form.get('phone')  # NEW: Get phone input
        password = request.form['password']
        security_q = request.form.get('security_q', '')
        security_a = request.form.get('security_a', '').lower().strip()

        with get_db_connection() as conn:
            try:
                # Update INSERT statement to include phone
                conn.execute('''
                    INSERT INTO users (username, email, phone, password, security_q, security_a) 
                    VALUES (?, ?, ?, ?, ?, ?)
                ''', (username, email, phone, generate_password_hash(password), security_q, security_a))
                
                conn.commit()
                flash('Account created! Please login.', 'success')
                return redirect('/login')
            except sqlite3.IntegrityError:
                flash('Username or email already exists.', 'danger')
                
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    """User login route."""
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        
        with get_db_connection() as conn:
            user = conn.execute('SELECT * FROM users WHERE username = ?', (username,)).fetchone()
        
        if user and check_password_hash(user['password'], password):
            session.clear()
            session['user_id'] = user['id']
            session['username'] = user['username']
            session['email'] = user['email']
            session['skill_level'] = user['skill_level']
            session['join_date'] = user['created_at'][:10]
            
            # Update session stats
            user_id = user['id']
            with get_db_connection() as conn:
                total_quizzes = conn.execute('SELECT COUNT(*) FROM quizzes WHERE user_id = ?', (user_id,)).fetchone()[0]
                avg_score_result = conn.execute('SELECT AVG(score) FROM quizzes WHERE user_id = ? AND status = ?', 
                                              (user_id, 'completed')).fetchone()[0]
            
            session['quizzes_taken'] = total_quizzes
            session['avg_score'] = round(avg_score_result, 1) if avg_score_result else 0
            
            flash('Login successful!', 'success')
            return redirect('/dashboard')
        else:
            flash('Invalid username or password', 'danger')
    
    return render_template('login.html')

@app.route('/forgot_password', methods=['GET', 'POST'])
def forgot_password():
    """Password recovery route."""
    if request.method == 'POST':
        stage = request.form.get('stage')
        username = request.form.get('username')
        
        with get_db_connection() as conn:
            user = conn.execute('SELECT * FROM users WHERE username = ?', (username,)).fetchone()
            
            if not user:
                flash('Username not found.', 'danger')
                if stage in ['verify_q', 'reset_p']:
                    return render_template('forgot_password.html', stage=stage)
                return render_template('forgot_password.html', stage='ask_username')

            if stage == 'verify_q':
                security_a = request.form.get('security_a', '').lower().strip()
                user_dict = dict(user)
                
                if security_a == user_dict['security_a'].lower().strip():
                    flash('Security question answered successfully. Reset your password now.', 'success')
                    return render_template('forgot_password.html', user=user_dict, stage='reset_p')
                else:
                    flash('Incorrect security answer.', 'danger')
                    return render_template('forgot_password.html', user=user_dict, stage='verify_q')

            elif stage == 'reset_p':
                new_password = request.form.get('new_password')
                confirm_password = request.form.get('confirm_password')

                if new_password != confirm_password:
                    flash('Passwords do not match.', 'danger')
                    user_dict = dict(user)
                    return render_template('forgot_password.html', user=user_dict, stage='reset_p')
                
                hashed_password = generate_password_hash(new_password)
                conn.execute('UPDATE users SET password = ? WHERE id = ?', (hashed_password, user['id']))
                conn.commit()
                
                flash('Password successfully reset! Please login.', 'success')
                return redirect(url_for('login'))
        
        if username:
            user_dict = dict(user)
            return render_template('forgot_password.html', user=user_dict, stage='verify_q')
        
        return render_template('forgot_password.html', stage='ask_username')
        
    return render_template('forgot_password.html', stage='ask_username')

@app.route('/dashboard')
def dashboard():
    """User dashboard route."""
    if not is_logged_in():
        return redirect('/login')
    
    user_id = session['user_id']
    with get_db_connection() as conn:
        total_quizzes = conn.execute('SELECT COUNT(*) FROM quizzes WHERE user_id = ?', (user_id,)).fetchone()[0]
        completed_quizzes = conn.execute('SELECT COUNT(*) FROM quizzes WHERE user_id = ? AND status = ?', 
                                   (user_id, 'completed')).fetchone()[0]
        recent_quizzes = conn.execute('SELECT * FROM quizzes WHERE user_id = ? ORDER BY created_at DESC LIMIT 5', 
                                    (user_id,)).fetchall()
    
    recent_quizzes_list = [dict(q) for q in recent_quizzes]
    
    # Add motivational quote to dashboard
    motivational_quote = random.choice(MOTIVATIONAL_QUOTES)

    return render_template('dashboard.html',
                         username=session['username'],
                         total_quizzes=total_quizzes,
                         completed_quizzes=completed_quizzes,
                         recent_quizzes=recent_quizzes_list,
                         motivational_quote=motivational_quote)

@app.route('/profile')
def profile():
    """User profile route."""
    if not is_logged_in():
        return redirect('/login')

    user = get_current_user_data(session['user_id'])
    
    user_data_display = dict(user)
    if user_data_display and user_data_display['created_at'] and isinstance(user_data_display['created_at'], str):
        try:
            user_data_display['created_at'] = datetime.strptime(user_data_display['created_at'].split('.')[0], '%Y-%m-%d %H:%M:%S')
        except ValueError:
            pass
    
    # Fetch performance analytics
    with get_db_connection() as conn:
        quizzes = conn.execute('SELECT id, topic, score, questions FROM quizzes WHERE user_id = ? AND status = ?', 
                               (user['id'], 'completed')).fetchall()

    topic_summary = {}
    total_completed = 0
    total_average_score = 0
    
    for q in quizzes:
        total_completed += 1
        total_average_score += q['score']
        
        topic = q['topic']
        try:
            questions = json.loads(q['questions'])
        except json.JSONDecodeError:
            continue
        
        answered_questions_in_quiz = [qn for qn in questions if 'user_answer' in qn]
        
        total_time_in_quiz = sum(qn.get('response_time', 0) for qn in answered_questions_in_quiz)
        total_q_in_quiz = len(answered_questions_in_quiz)
        correct_q_in_quiz = sum(1 for qn in answered_questions_in_quiz if qn.get('is_correct', False))
        
        if topic not in topic_summary:
            topic_summary[topic] = {
                'total_score': 0, 
                'quiz_count': 0,
                'total_response_time': 0,
                'total_questions': 0,
                'total_correct_answers': 0 
            }
        
        summary = topic_summary[topic]
        summary['total_score'] += q['score']
        summary['quiz_count'] += 1
        summary['total_response_time'] += total_time_in_quiz
        summary['total_questions'] += total_q_in_quiz
        summary['total_correct_answers'] += correct_q_in_quiz

    performances_list = []
    for topic, summary in topic_summary.items():
        avg_score = summary['total_score'] / summary['quiz_count'] if summary['quiz_count'] > 0 else 0
        avg_accuracy = (summary['total_correct_answers'] / summary['total_questions']) if summary['total_questions'] > 0 else 0
        avg_time = summary['total_response_time'] / summary['total_questions'] if summary['total_questions'] > 0 else 0
        
        performances_list.append({
            'topic': topic,
            'difficulty': user['skill_level'],
            'accuracy': avg_accuracy,      
            'total_questions': summary['total_questions'],
            'correct_answers': summary['total_correct_answers'], 
            'average_response_time': avg_time
        })
    
    overall_avg_score = total_average_score / total_completed if total_completed > 0 else 0
    
    # Add motivational quote
    motivational_quote = random.choice(MOTIVATIONAL_QUOTES)
    
    return render_template('profile.html', 
                           user=user_data_display,
                           total_quizzes=total_completed,
                           completed_quizzes=total_completed,
                           average_score=overall_avg_score,
                           performances=performances_list,
                           motivational_quote=motivational_quote)

@app.route('/update_profile', methods=['POST'])
def update_profile():
    """Update user profile route."""
    if not is_logged_in():
        return redirect('/login')
    
    user_id = session['user_id']
    username_new = request.form.get('username', '').strip()
    skill_level_new = request.form.get('skill_level', 'Medium')
    
    with get_db_connection() as conn:
        existing_user = conn.execute('SELECT id FROM users WHERE username = ? AND id != ?', 
                                   (username_new, user_id)).fetchone()
        
        if existing_user:
            flash('Username is already taken.', 'danger')
            return redirect('/profile')
            
        conn.execute('''
            UPDATE users 
            SET username = ?, 
                skill_level = ?
            WHERE id = ?
        ''', (username_new, skill_level_new, user_id))
        
        conn.commit()
        
        session['username'] = username_new
        session['skill_level'] = skill_level_new
    
    flash('Profile updated successfully!', 'success')
    return redirect('/profile')

# =========================================================================
# === QUIZ CREATION AND SUBMISSION ROUTES
# =========================================================================

# app.py - CRITICAL SECTION: /create_quiz route

@app.route('/create_quiz', methods=['GET', 'POST'])
def create_quiz():
    """Create new quiz with URL Scraping Support."""
    if not is_logged_in():
        return redirect('/login')

    user_id = session['user_id']
    current_difficulty = get_user_skill_level(user_id)
    
    if request.method == 'POST':
        topic = request.form.get('topic', 'General').strip()
        user_difficulty = request.form.get('difficulty', current_difficulty).capitalize()
        num_questions = int(request.form.get('num_questions', 5))
        quiz_type = request.form.get('quiz_type', 'adaptive')
        
        # 1. Get input
        raw_input = request.form.get('content', '').strip()
        final_content = raw_input
        
        # 2. Check if it looks like a URL
        if raw_input.startswith(('http://', 'https://')):
            print(f"DEBUG: Detected URL: {raw_input}")
            scraped_text = scrape_url_content(raw_input)
            if scraped_text and len(scraped_text) > 100:
                final_content = scraped_text
                print("DEBUG: Successfully scraped content.")
            else:
                flash("Could not read content from that URL. Using topic instead.", "warning")
                final_content = f"Generate questions about {topic}"

        # 3. Fallback
        if not final_content:
            final_content = f"Generate questions about {topic}"

        # 4. Generate Questions
        # We request 3x the questions for adaptive mode to ensure we have Hard/Easy options
        pool_size = num_questions * 3 if quiz_type == 'adaptive' else num_questions
        
        question_pool = quiz_engine.generate_questions(
            content=final_content, 
            num_questions=pool_size,
            topic=topic
        )
        
        if not question_pool:
            flash('Could not generate quiz questions. Please try again.', 'danger')
            return render_template('create_quiz.html', current_difficulty=current_difficulty)

        quiz_title = f"{topic} Quiz ({user_difficulty} - {quiz_type.capitalize()})"
        
        with get_db_connection(dict_cursor=False) as conn:
            cursor = conn.cursor()
            cursor.execute('''
                INSERT INTO quizzes (user_id, title, topic, content, questions, difficulty, quiz_type)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            ''', (user_id, quiz_title, topic, final_content, json.dumps(question_pool), user_difficulty, quiz_type))
            
            quiz_id = cursor.lastrowid
            conn.commit()
            
        flash(f'Quiz "{quiz_title}" created successfully!', 'success')
        
        session['user_quiz_length'] = num_questions 
        
        if quiz_type == 'adaptive':
            return redirect(url_for('take_quiz', quiz_id=quiz_id))
        else:
            return redirect(url_for('take_simple_quiz', quiz_id=quiz_id))

    return render_template('create_quiz.html', current_difficulty=current_difficulty)
# ... (Rest of app.py is unchanged)

@app.route('/quiz/<int:quiz_id>')
def take_quiz(quiz_id):
    """Enhanced adaptive quiz route with graceful initial question selection."""
    if not is_logged_in():
        return redirect('/login')
    
    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ? AND user_id = ? AND quiz_type = ?', 
                           (quiz_id, session['user_id'], 'adaptive')).fetchone()
    
    if not quiz:
        flash('Adaptive Quiz not found or is a simple quiz.', 'danger')
        return redirect('/dashboard')
    
    questions_pool = json.loads(quiz['questions'])

    # CRITICAL CHECK: Ensure the pool is not empty before proceeding
    if not questions_pool:
        flash('Quiz question pool is empty. Please try recreating the quiz with different content.', 'danger')
        return redirect('/dashboard')
    
    # Initialize session tracking for the quiz
    if session.get('current_quiz_id') != quiz_id:
        session['current_quiz_id'] = quiz_id
        session['current_question_index'] = 0
        session['score'] = 0
        session['answered_questions'] = []
        
        start_difficulty = quiz['difficulty'].lower()
        if hasattr(adaptive_engine, 'get_difficulty_index'):
             session['current_difficulty_index'] = adaptive_engine.get_difficulty_index(start_difficulty)
        else:
             session['current_difficulty_index'] = 1 # Default to medium index
        
        initial_q = None
        
        # 1. Try to find question matching starting difficulty
        initial_q = next((q for q in questions_pool if q.get('difficulty', 'medium').lower() == start_difficulty), None)
        
        # 2. Fallback using difficulty throttling logic
        if not initial_q:
            fallback_order = ['medium', 'easy', 'hard']
            
            for fallback_diff in fallback_order:
                # Find an unanswered question matching the fallback difficulty
                initial_q = next((q for q in questions_pool if q.get('difficulty', 'medium').lower() == fallback_diff and q['id'] not in {q['q_id'] for q in session.get('answered_questions', [])}), None)
                if initial_q:
                    break

        if initial_q:
            quiz_length = session.get('user_quiz_length', 10) 
            return render_template('single_question_quiz.html', 
                                 quiz=dict(quiz), 
                                 questions=[initial_q],
                                 quiz_length=quiz_length,
                                 current_index=session['current_question_index'], 
                                 show_evaluation=False)
        else:
            flash('No starting question could be found. Question pool is invalid.', 'danger')
            return redirect('/dashboard')
    
    return redirect(url_for('next_question_adaptive', quiz_id=quiz_id))

@app.route('/submit_answer/<int:quiz_id>', methods=['POST'])
def submit_answer(quiz_id):
    """Enhanced answer submission with multiple question type support."""
    if not is_logged_in() or session.get('current_quiz_id') != quiz_id:
        return redirect('/dashboard')
    
    current_difficulty_index = session.get('current_difficulty_index', 1)
    
    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ? AND user_id = ?', 
                           (quiz_id, session['user_id'])).fetchone()
    
    if not quiz:
        flash('Quiz not found', 'danger')
        return redirect('/dashboard')
    
    # Get submitted data
    question_id = int(request.form.get('question_id', 0))
    user_answer = request.form.get('answer', '').strip()
    time_taken = int(request.form.get('time_taken', 1))
    
    # Handle checkbox answers (multiple selections)
    if request.form.getlist('answer'):  # For checkbox questions
        user_answer = ','.join(request.form.getlist('answer'))
    
    # Find current question
    questions_pool = json.loads(quiz['questions'])
    current_question = next((q for q in questions_pool if q.get('id') == question_id), None)
    
    if not current_question:
        flash("Could not locate the current question in the pool. Quitting.", 'danger')
        return redirect(url_for('finalize_quiz', quiz_id=quiz_id))

    # Evaluate answer using enhanced evaluation function
    is_correct = evaluate_answer(current_question, user_answer)
    
    # Handle feedback
    feedback_type = request.form.get('feedback_type', '').strip().lower()
    feedback_comment = request.form.get('feedback_comment', '').strip()
    flag_question = 1 if request.form.get('flag_question') else 0

    if feedback_type not in {'positive', 'negative', 'neutral'}:
        feedback_type = None

    if feedback_type or feedback_comment or flag_question:
        with get_db_connection(dict_cursor=False) as feedback_conn:
            feedback_conn.execute(
                '''
                INSERT INTO question_feedback (quiz_id, question_id, user_id, feedback_type, comment, flagged, question_text)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                ''',
                (
                    quiz_id,
                    str(question_id),
                    session['user_id'],
                    feedback_type,
                    feedback_comment,
                    flag_question,
                    current_question.get('question_text', '')
                )
            )
            feedback_conn.commit()

    # Update session data
    if is_correct:
        session['score'] = session.get('score', 0) + 1
        
    session['answered_questions'].append({
        'q_id': question_id,
        'user_answer': user_answer, 
        'difficulty': current_question['difficulty'],
        'is_correct': is_correct,
        'response_time': time_taken,
        'question_type': current_question.get('question_type', 'mcq')
    })
    
    # Adaptive difficulty adjustment
    if is_correct:
        current_difficulty_index = min(current_difficulty_index + 1, 2)
        flash('✅ Correct! Increasing difficulty for the next question.', 'success')
    else:
        current_difficulty_index = max(current_difficulty_index - 1, 0)
        flash('❌ Incorrect. Decreasing difficulty for the next question.', 'warning')
    
    session['current_difficulty_index'] = current_difficulty_index
    session['current_question_index'] = session.get('current_question_index', 0) + 1
    
    # Update question with response details
    current_question['user_answer'] = user_answer
    current_question['is_correct'] = is_correct
    current_question['response_time'] = time_taken
    
    quiz_length = session.get('user_quiz_length', 10)

    # Show evaluation
    return render_template('single_question_quiz.html',
                         quiz=dict(quiz),
                         questions=[current_question], 
                         quiz_length=quiz_length,
                         current_index=session['current_question_index'],
                         show_evaluation=True,
                         user_answer=user_answer,
                         is_correct=is_correct,
                         correct_answer=current_question['correct_answer'],
                         explanation=current_question.get('explanation', 'No explanation available.'))

@app.route('/next_question_adaptive/<int:quiz_id>')
def next_question_adaptive(quiz_id):
    """Find and display next adaptive question with difficulty throttling."""
    if not is_logged_in() or session.get('current_quiz_id') != quiz_id:
        return redirect('/dashboard')
    
    answered_count = session['current_question_index']
    quiz_length = session.get('user_quiz_length', 10)
    
    # Stop after requested number of questions
    if answered_count >= quiz_length: 
        return redirect(url_for('finalize_quiz', quiz_id=quiz_id))
    
    # Determine required difficulty
    required_difficulty_index = session.get('current_difficulty_index', 1)
    
    if hasattr(adaptive_engine, 'get_difficulty_by_index'):
        required_difficulty = adaptive_engine.get_difficulty_by_index(required_difficulty_index).lower()
    else:
        required_difficulty = 'medium'


    # Find next question
    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ?', (quiz_id,)).fetchone()
        questions_pool = json.loads(quiz['questions'])
    
    answered_q_ids = {q['q_id'] for q in session.get('answered_questions', [])}
    
    next_question = None
    
    # MODIFIED LOGIC: Implement difficulty fall-back (Throttling)
    # 1. Try required difficulty
    required_order = [required_difficulty]
    
    # 2. Add adjacent difficulties as fallback
    if required_difficulty == 'hard':
        required_order.extend(['medium', 'easy'])
    elif required_difficulty == 'easy':
        required_order.extend(['medium', 'hard'])
    else: # medium
        if random.choice([True, False]):
            required_order.extend(['hard', 'easy'])
        else:
            required_order.extend(['easy', 'hard'])

    # Search through the prioritized order of difficulties
    for difficulty in required_order:
        next_question = next((q for q in questions_pool 
                               if q.get('difficulty', 'medium').lower() == difficulty and q['id'] not in answered_q_ids), None)
        if next_question:
            if difficulty != required_difficulty:
                # Notify the user that throttling occurred
                flash(f"⚠️ Could not find a **{required_difficulty.capitalize()}** question. Using a **{difficulty.capitalize()}** question instead to keep the quiz flowing.", 'info')
            break
            
    if next_question:
        return render_template('single_question_quiz.html',
                             quiz=dict(quiz),
                             questions=[next_question],
                             quiz_length=quiz_length,
                             current_index=answered_count, 
                             show_evaluation=False)
    else:
        # Final fallback: just get any remaining un-answered question, regardless of difficulty
        next_question_any = next((q for q in questions_pool if q['id'] not in answered_q_ids), None)
        
        if next_question_any:
            flash(f"⚠️ Question pool exhausted for the adaptive path. Using the next available question.", 'info')
            return render_template('single_question_quiz.html',
                                 quiz=dict(quiz),
                                 questions=[next_question_any],
                                 quiz_length=quiz_length,
                                 current_index=answered_count, 
                                 show_evaluation=False)
        else:
            flash("You've exhausted the question pool. Finalizing quiz.", 'info')
            return redirect(url_for('finalize_quiz', quiz_id=quiz_id))

@app.route('/finalize_quiz/<int:quiz_id>')
def finalize_quiz(quiz_id):
    """Finalize adaptive quiz and save results."""
    if not is_logged_in() or session.get('current_quiz_id') != quiz_id:
        return redirect('/dashboard')

    answered_log = session.get('answered_questions', [])
    total_answered = len(answered_log)
    total_correct = session.get('score', 0)
    
    final_score = (total_correct / total_answered) * 100 if total_answered > 0 else 0
    
    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ?', (quiz_id,)).fetchone()
        
        # Update Performance Table
        full_pool = json.loads(quiz['questions'])
        final_questions_map = {q['id']: q for q in full_pool}

        for log in answered_log:
            q_detail = final_questions_map.get(log['q_id'], {})
            
            adaptive_engine.update_performance(
                user_id=session['user_id'],
                topic=quiz['topic'],
                difficulty=q_detail.get('difficulty', 'medium'),
                is_correct=log['is_correct'],
                response_time=log['response_time']
            )
            
            # Merge user response data
            if log['q_id'] in final_questions_map:
                final_questions_map[log['q_id']]['user_answer'] = log.get('user_answer', '') 
                final_questions_map[log['q_id']]['is_correct'] = log['is_correct']
                final_questions_map[log['q_id']]['response_time'] = log['response_time']

        # Update user skill level
        next_difficulty = adaptive_engine.calculate_next_difficulty(
            user_id=session['user_id'],
            topic=quiz['topic'],
            current_score=final_score
        )
        conn.execute('UPDATE users SET skill_level = ? WHERE id = ?', 
                     (next_difficulty.capitalize(), session['user_id']))
        
        session['skill_level'] = next_difficulty.capitalize()

        # Update quiz status
        conn.execute('''UPDATE quizzes 
                       SET score = ?, status = ?, questions = ?
                       WHERE id = ?''',
                    (final_score, 'completed', json.dumps(list(final_questions_map.values())), quiz_id))
        conn.commit()
    
    # Clear session
    session.pop('current_quiz_id', None)
    session.pop('current_question_index', None)
    session.pop('score', None)
    session.pop('answered_questions', None)
    session.pop('current_difficulty_index', None)
    session.pop('user_quiz_length', None)
    
    # Update session stats
    user_id = session['user_id']
    with get_db_connection() as conn:
        total_quizzes = conn.execute('SELECT COUNT(*) FROM quizzes WHERE user_id = ?', (user_id,)).fetchone()[0]
        avg_score_result = conn.execute('SELECT AVG(score) FROM quizzes WHERE user_id = ? AND status = ?', 
                                      (user_id, 'completed')).fetchone()[0]
    session['quizzes_taken'] = total_quizzes
    session['avg_score'] = round(avg_score_result, 1) if avg_score_result else 0
    
    return redirect(f'/performance/{quiz_id}')

# =========================================================================
# === SIMPLE QUIZ FLOW
# =========================================================================

@app.route('/quiz_simple/<int:quiz_id>')
def take_simple_quiz(quiz_id):
    """Simple quiz route (all questions on one page)."""
    if not is_logged_in():
        return redirect('/login')
    
    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ? AND user_id = ? AND quiz_type = ?', 
                           (quiz_id, session['user_id'], 'simple')).fetchone()
    
    if not quiz:
        flash('Simple Quiz not found or is an adaptive quiz.', 'danger')
        return redirect('/dashboard')
    
    if quiz['status'] == 'completed':
        return redirect(url_for('performance_analysis', quiz_id=quiz_id))
        
    questions = json.loads(quiz['questions'])
    
    return render_template('simple_quiz.html', quiz=dict(quiz), questions=questions)

@app.route('/submit_quiz/<int:quiz_id>', methods=['POST'])
def submit_quiz(quiz_id):
    """Submit simple quiz with multiple question type support."""
    if not is_logged_in():
        return redirect('/login')

    user_id = session['user_id']
    
    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ? AND user_id = ? AND quiz_type = ?', 
                           (quiz_id, user_id, 'simple')).fetchone()

    if not quiz or quiz['status'] == 'completed':
        flash('Quiz not found or already completed.', 'danger')
        return redirect('/dashboard')

    # Evaluate answers
    questions = json.loads(quiz['questions'])
    total_questions = len(questions)
    correct_count = 0
    answered_log = []
    feedback_entries = []
    
    question_map = {q['id']: q for q in questions}

    for question in questions:
        q_id = question['id']
        
        # Handle different question types
        if question.get('question_type') == 'checkbox':
            user_answer = ','.join(request.form.getlist(f'answer_{q_id}'))
        else:
            user_answer = request.form.get(f'answer_{q_id}', '').strip()
            
        time_taken = int(request.form.get(f'time_{q_id}', 1))
        
        # Evaluate answer
        is_correct = evaluate_answer(question, user_answer)
        
        if is_correct:
            correct_count += 1
            
        # Update question with user response
        question['user_answer'] = user_answer
        question['is_correct'] = is_correct
        question['response_time'] = time_taken
        
        # Prepare performance log
        answered_log.append({
            'q_id': q_id,
            'user_answer': user_answer,
            'difficulty': question['difficulty'],
            'is_correct': is_correct,
            'response_time': time_taken,
            'question_type': question.get('question_type', 'mcq')
        })

        # Collect feedback
        feedback_type = request.form.get(f'feedback_type_{q_id}', '').strip().lower()
        feedback_comment = request.form.get(f'feedback_comment_{q_id}', '').strip()
        flag_question = 1 if request.form.get(f'flag_question_{q_id}') else 0

        if feedback_type not in {'positive', 'negative', 'neutral'}:
            feedback_type = None

        if feedback_type or feedback_comment or flag_question:
            feedback_entries.append({
                'quiz_id': quiz_id,
                'question_id': str(q_id),
                'feedback_type': feedback_type,
                'feedback_comment': feedback_comment,
                'flagged': flag_question,
                'question_text': question.get('question_text', '')
            })
    
    final_score = (correct_count / total_questions) * 100 if total_questions > 0 else 0
    
    # Save feedback
    if feedback_entries:
        with get_db_connection(dict_cursor=False) as feedback_conn:
            for entry in feedback_entries:
                feedback_conn.execute(
                    '''
                    INSERT INTO question_feedback (quiz_id, question_id, user_id, feedback_type, comment, flagged, question_text)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                    ''',
                    (
                        entry['quiz_id'],
                        entry['question_id'],
                        user_id,
                        entry['feedback_type'],
                        entry['feedback_comment'],
                        entry['flagged'],
                        entry['question_text']
                    )
                )
            feedback_conn.commit()

    # Update performance and skill level
    for log in answered_log:
        q_detail = question_map.get(log['q_id'], {})
        adaptive_engine.update_performance(
            user_id=user_id,
            topic=quiz['topic'],
            difficulty=q_detail.get('difficulty', 'medium'),
            is_correct=log['is_correct'],
            response_time=log['response_time']
        )
    
    next_difficulty = adaptive_engine.calculate_next_difficulty(
        user_id=user_id,
        topic=quiz['topic'],
        current_score=final_score
    )
    conn.execute('UPDATE users SET skill_level = ? WHERE id = ?', 
                 (next_difficulty.capitalize(), user_id))
    session['skill_level'] = next_difficulty.capitalize()

    # Update quiz
    conn.execute('''UPDATE quizzes 
                   SET score = ?, status = ?, questions = ?
                   WHERE id = ?''',
                (final_score, 'completed', json.dumps(questions), quiz_id))
    conn.commit()

    # Update session stats
    with get_db_connection() as conn_stats:
        total_quizzes_stat = conn_stats.execute('SELECT COUNT(*) FROM quizzes WHERE user_id = ?', (user_id,)).fetchone()[0]
        avg_score_result = conn_stats.execute('SELECT AVG(score) FROM quizzes WHERE user_id = ? AND status = ?', 
                                      (user_id, 'completed')).fetchone()[0]
    session['quizzes_taken'] = total_quizzes_stat
    session['avg_score'] = round(avg_score_result, 1) if avg_score_result else 0
    
    flash(f'Quiz completed! Your score: {final_score:.1f}%', 'success')
    return redirect(url_for('performance_analysis', quiz_id=quiz_id))

@app.route('/logout')
def logout():
    """Logout route."""
    session.clear()
    flash('You have been logged out.', 'info')
    return redirect('/')

@app.route('/performance/<int:quiz_id>')
def performance_analysis(quiz_id):
    """Enhanced performance analysis with AI suggestions link."""
    if not is_logged_in():
        return redirect('/login')

    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ? AND user_id = ?', 
                           (quiz_id, session['user_id'])).fetchone()

    if not quiz or quiz['status'] != 'completed':
        flash('Quiz not found or not completed.', 'danger')
        return redirect('/dashboard')
    
    questions = json.loads(quiz['questions'])
    
    # Filter answered questions
    final_questions = [q for q in questions if 'user_answer' in q]

    total_questions = len(final_questions)
    correct_count = sum(1 for q in final_questions if q.get('is_correct'))
    
    score = (correct_count / total_questions) * 100 if total_questions else 0
    incorrect_count = total_questions - correct_count

    # Enhanced breakdowns
    difficulty_breakdown = {'Easy': {'correct': 0, 'total': 0}, 'Medium': {'correct': 0, 'total': 0}, 'Hard': {'correct': 0, 'total': 0}}
    type_breakdown = {
        'mcq': {'correct': 0, 'total': 0}, 
        'checkbox': {'correct': 0, 'total': 0}, 
        'true_false': {'correct': 0, 'total': 0}, 
        'short_answer': {'correct': 0, 'total': 0},
        'dropdown': {'correct': 0, 'total': 0}
    }
    total_time = 0
    time_per_q_list = []
    
    for q in final_questions:
        # --- FIX: Ensure response_time exists to prevent Template Error ---
        if 'response_time' not in q or q['response_time'] is None:
            q['response_time'] = 0
        # ------------------------------------------------------------------

        diff = q.get('difficulty', 'Medium').capitalize()
        q_type = q.get('question_type', 'mcq').lower()
        is_correct = q.get('is_correct', False)
        response_time = q.get('response_time', 0)
        
        # Update difficulty breakdown
        difficulty_breakdown.setdefault(diff, {'correct': 0, 'total': 0})['total'] += 1
        if is_correct:
            difficulty_breakdown[diff]['correct'] += 1
        
        # Update type breakdown
        type_breakdown.setdefault(q_type, {'correct': 0, 'total': 0})['total'] += 1
        if is_correct:
            type_breakdown[q_type]['correct'] += 1
            
        total_time += response_time
        time_per_q_list.append(response_time)

    avg_time_per_q = total_time / total_questions if total_questions else 0

    # Select motivational quote for performance page too
    motivational_quote = random.choice(MOTIVATIONAL_QUOTES)

    return render_template(
        'performance_analysis.html', 
        quiz=dict(quiz),
        score=score, 
        correct_count=correct_count,
        incorrect_count=incorrect_count,
        total_questions=total_questions,
        difficulty_breakdown=difficulty_breakdown,
        type_breakdown=type_breakdown,
        total_time=total_time,
        avg_time_per_q=avg_time_per_q,
        questions=final_questions,
        time_per_q_list=time_per_q_list,
        motivational_quote=motivational_quote,
        quiz_id=quiz_id
    )

@app.route('/ai_suggestions/<int:quiz_id>')
def ai_suggestions(quiz_id):
    """Generate AI-powered learning suggestions based on quiz performance."""
    if not is_logged_in():
        return redirect('/login')

    with get_db_connection() as conn:
        quiz = conn.execute('SELECT * FROM quizzes WHERE id = ? AND user_id = ?', 
                           (quiz_id, session['user_id'])).fetchone()

    if not quiz or quiz['status'] != 'completed':
        flash('Quiz not found or not completed.', 'danger')
        return redirect(f'/performance/{quiz_id}')
    
    questions = json.loads(quiz['questions'])
    answered_questions = [q for q in questions if 'user_answer' in q]
    
    if not answered_questions:
        flash('No questions found for analysis.', 'danger')
        return redirect(f'/performance/{quiz_id}')
    
    # Analyze performance
    total_questions = len(answered_questions)
    correct_count = sum(1 for q in answered_questions if q.get('is_correct'))
    score = (correct_count / total_questions) * 100
    
    # Identify weak areas
    incorrect_questions = [q for q in answered_questions if not q.get('is_correct')]
    weak_topics = {}
    
    for q in incorrect_questions:
        topic = q.get('topic', 'General')
        weak_topics[topic] = weak_topics.get(topic, 0) + 1
    
    # Difficulty analysis
    difficulty_performance = {}
    for q in answered_questions:
        diff = q.get('difficulty', 'Medium')
        if diff not in difficulty_performance:
            difficulty_performance[diff] = {'correct': 0, 'total': 0}
        difficulty_performance[diff]['total'] += 1
        if q.get('is_correct'):
            difficulty_performance[diff]['correct'] += 1
    
    # Generate AI suggestions
    suggestions = generate_ai_suggestions(score, weak_topics, difficulty_performance, quiz['topic'])
    
    # Select random motivational quote
    motivational_quote = random.choice(MOTIVATIONAL_QUOTES)
    
    # Voice message based on performance
    if score >= 90:
        voice_message = f"Outstanding! You scored {score:.1f}%! You're mastering {quiz['topic']}!"
        performance_level = "excellent"
    elif score >= 75:
        voice_message = f"Great work! {score:.1f}% shows solid understanding of {quiz['topic']}!"
        performance_level = "good"
    elif score >= 60:
        voice_message = f"Good effort! {score:.1f}% - keep practicing to improve your {quiz['topic']} skills!"
        performance_level = "average"
    else:
        voice_message = f"Don't give up! {score:.1f}% is a starting point. Focus on the basics of {quiz['topic']}!"
        performance_level = "improving"
    
    return render_template('ai_suggestions.html',
                         quiz=dict(quiz),
                         score=score,
                         correct_count=correct_count,
                         total_questions=total_questions,
                         suggestions=suggestions,
                         weak_topics=weak_topics,
                         motivational_quote=motivational_quote,
                         voice_message=voice_message,
                         performance_level=performance_level)

@app.route('/results/<int:quiz_id>')
def quiz_results(quiz_id):
    """Alias for performance analysis."""
    return redirect(f'/performance/{quiz_id}')

@app.route('/leaderboard/<int:quiz_id>')
def leaderboard(quiz_id):
    """Leaderboard route showing top scores for the quiz topic."""
    if not is_logged_in():
        return redirect('/login')
    
    with get_db_connection() as conn:
        # 1. Get the details of the quiz to identify Topic and Difficulty
        current_quiz = conn.execute('SELECT * FROM quizzes WHERE id = ?', (quiz_id,)).fetchone()
        
        if not current_quiz:
            flash('Quiz not found.', 'danger')
            return redirect('/dashboard')

        topic = current_quiz['topic']
        difficulty = current_quiz['difficulty']

        # 2. Get top 50 attempts for this Topic + Difficulty
        # We join with the users table to get usernames
        leaderboard_data = conn.execute('''
            SELECT u.username, u.id as user_id, q.score, q.created_at, q.questions
            FROM quizzes q
            JOIN users u ON q.user_id = u.id
            WHERE q.topic = ? 
              AND q.difficulty = ? 
              AND q.status = 'completed'
            ORDER BY q.score DESC
            LIMIT 50
        ''', (topic, difficulty)).fetchall()

    # 3. Process data (Calculate total time from JSON data since it's not a column)
    attempts = []
    for row in leaderboard_data:
        total_time = 0
        try:
            # Parse questions JSON to sum up response times
            questions_data = json.loads(row['questions'])
            total_time = sum(q.get('response_time', 0) for q in questions_data if 'user_answer' in q)
        except:
            total_time = 0

        # Handle date parsing safely
        completed_at = row['created_at']
        if isinstance(completed_at, str):
            try:
                completed_at = datetime.strptime(completed_at.split('.')[0], '%Y-%m-%d %H:%M:%S')
            except ValueError:
                completed_at = datetime.now() # Fallback

        attempts.append({
            'username': row['username'],
            'user_id': row['user_id'],
            'score': row['score'],
            'total_time': total_time,
            'completed_at': completed_at
        })

    # 4. Sort: Primary = Score (High to Low), Secondary = Time (Low to High)
    attempts.sort(key=lambda x: (-x['score'], x['total_time']))

    return render_template('leaderboard.html', quiz=dict(current_quiz), attempts=attempts)

@app.route('/feedback', methods=['GET', 'POST'])
def feedback():
    """User feedback route."""
    if not is_logged_in():
        flash('You must be logged in to leave feedback.', 'danger')
        return redirect('/login')
        
    if request.method == 'POST':
        feedback_type = request.form.get('feedback_type')
        rating = request.form.get('rating')
        comments = request.form.get('comments', '').strip()
        
        user_id = session['user_id']
        username = session['username']
        
        # Log feedback (in a real app, save to database)
        print("-" * 40)
        print(f"🚨 NEW FEEDBACK RECEIVED (User ID: {user_id}, Username: {username})")
        print(f"Type: {feedback_type}")
        print(f"Rating: {rating}/5")
        print(f"Comments: {comments}")
        print("-" * 40)
        
        flash(f'Thank you for your feedback ({feedback_type})! We appreciate you helping us improve.', 'success')
        return redirect('/dashboard')
        
    return render_template('feedback.html')

if __name__ == '__main__':
    print("🚀 SMART QUIZZER STARTING...")
    print("✅ Database initialized")
    print("✅ AI Engine ready")
    print("✅ Flask app configured")
    print("✅ Motivational features loaded")
    print("✅ Voice synthesis ready")
    
    with app.app_context():
        init_db()
    
    # Run on a common port for web applications
    app.run(debug=True, port=5000)
